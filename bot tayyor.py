from aiogram import Bot, Dispatcher, types
from aiogram.types import ReplyKeyboardMarkup, KeyboardButton, ReplyKeyboardRemove
from aiogram.utils import executor
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters.state import State, StatesGroup

from docx import Document
from pptx import Presentation
import os, json
from datetime import datetime, timedelta
import google.generativeai as genai

API_TOKEN = os.getenv("BOT_TOKEN")
GEMINI_KEY = os.getenv("GEMINI_KEY")
ADMIN_ID = int(os.getenv("ADMIN_ID"))

LIMIT = 5

genai.configure(api_key=GEMINI_KEY)
model = genai.GenerativeModel("gemini-pro")

bot = Bot(token=API_TOKEN)
storage = MemoryStorage()
dp = Dispatcher(bot, storage=storage)

users_file = "users.json"

# ================= USERS =================
if os.path.exists(users_file):
    with open(users_file, "r") as f:
        users_data = json.load(f)
else:
    users_data = {}

def save_users():
    with open(users_file, "w") as f:
        json.dump(users_data, f)

def can_use(user_id):
    now = datetime.now().timestamp()
    user = users_data.get(str(user_id))

    if not user:
        users_data[str(user_id)] = {"count": 0, "timestamp": now}
        save_users()
        return True

    last = datetime.fromtimestamp(user["timestamp"])
    if datetime.now() - last > timedelta(days=30):
        user["count"] = 0
        user["timestamp"] = now
        save_users()

    return user["count"] < LIMIT

def add_use(user_id):
    users_data[str(user_id)]["count"] += 1
    save_users()

# ================= FSM =================
class BotStates(StatesGroup):
    choosing_degree = State()
    choosing_task = State()
    waiting_for_topic = State()

# ================= CHIROYLI BUTTONLAR =================
degree_keyboard = ReplyKeyboardMarkup(resize_keyboard=True)
degree_keyboard.row(
    KeyboardButton("🏫 Maktab"),
    KeyboardButton("🎓 Texnikum")
)
degree_keyboard.row(
    KeyboardButton("👩‍🏫 Universitet")
)

task_keyboard = ReplyKeyboardMarkup(resize_keyboard=True)
task_keyboard.row(
    KeyboardButton("📚 Dars ishlanma"),
    KeyboardButton("📝 Tezis")
)
task_keyboard.row(
    KeyboardButton("📄 Maqola"),
    KeyboardButton("🧪 Test")
)
task_keyboard.row(
    KeyboardButton("📊 Taqdimot")
)

# ================= START =================
@dp.message_handler(commands=['start'], state='*')
async def start(message: types.Message, state: FSMContext):
    await state.finish()
    await message.answer("👋 Salom! Darajani tanlang:", reply_markup=degree_keyboard)
    await BotStates.choosing_degree.set()

# ================= DEGREE =================
@dp.message_handler(lambda m: m.text in ["🏫 Maktab","🎓 Texnikum","👩‍🏫 Universitet"], state=BotStates.choosing_degree)
async def degree_selected(message: types.Message, state: FSMContext):
    await state.update_data(degree=message.text)
    await message.answer("📌 Vazifani tanlang:", reply_markup=task_keyboard)
    await BotStates.choosing_task.set()

# ================= TASK =================
@dp.message_handler(lambda m: m.text in [
    "📚 Dars ishlanma","📝 Tezis","📄 Maqola","🧪 Test","📊 Taqdimot"
], state=BotStates.choosing_task)
async def task_selected(message: types.Message, state: FSMContext):
    await state.update_data(task=message.text)
    await message.answer("✏️ Mavzuni yozing:", reply_markup=ReplyKeyboardRemove())
    await BotStates.waiting_for_topic.set()

# ================= GEMINI AI =================
def generate_ai_text(degree, task, topic):
    prompt = f"""
    {degree} uchun {task} tayyorla.
    Mavzu: {topic}
    Tuzilishi to'liq va professional bo'lsin.
    """

    response = model.generate_content(prompt)
    return response.text

# ================= TOPIC =================
@dp.message_handler(state=BotStates.waiting_for_topic)
async def topic_received(message: types.Message, state: FSMContext):
    user_id = message.from_user.id

    if not can_use(user_id):
        await message.answer("❌ Limit tugadi. Admin bilan bog'laning.")
        await state.finish()
        return

    data = await state.get_data()
    degree = data["degree"]
    task = data["task"]
    topic = message.text

    add_use(user_id)

    try:
        text = generate_ai_text(degree, task, topic)
    except:
        await message.answer("⚠️ AI javob bermadi. Keyinroq urinib ko‘ring.")
        await state.finish()
        return

    if task == "📊 Taqdimot":
        filename = f"{user_id}.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = topic
        slide.placeholders[1].text = text[:4000]
        prs.save(filename)
    else:
        filename = f"{user_id}.docx"
        doc = Document()
        doc.add_heading(task, 0)
        doc.add_paragraph(text)
        doc.save(filename)

    await message.answer_document(open(filename, "rb"), caption="✅ Tayyor")
    os.remove(filename)

    await state.finish()
    await message.answer("Yana boshlaymiz 👇", reply_markup=degree_keyboard)
    await BotStates.choosing_degree.set()

# ================= ADMIN PANEL =================
@dp.message_handler(commands=['admin'])
async def admin_panel(message: types.Message):
    if message.from_user.id != ADMIN_ID:
        return

    text = "👑 Admin panel\n\nFoydalanuvchilar:\n"
    for uid, info in users_data.items():
        text += f"{uid} → {info['count']}/{LIMIT}\n"

    await message.answer(text)

# ================= RUN =================
if __name__ == "__main__":
    print("Bot ishladi 🚀")
    executor.start_polling(dp, skip_updates=True)

